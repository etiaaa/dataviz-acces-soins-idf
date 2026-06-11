# -*- coding: utf-8 -*-
"""Génère le support de présentation (.pptx) du projet Accès aux soins IDF."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----- charte -----
BG      = RGBColor(0x1A, 0x1A, 0x1A)   # fond sombre (charte Think Positif)
WHITE   = RGBColor(0xF2, 0xF2, 0xF2)
GREY    = RGBColor(0xA8, 0xA8, 0xA8)
ACCENT  = RGBColor(0x4F, 0xC3, 0xF7)   # bleu clair = accent
GREEN   = RGBColor(0x66, 0xBB, 0x6A)   # bien doté
RED     = RGBColor(0xE5, 0x73, 0x73)   # tendu
ORANGE  = RGBColor(0xFF, 0xB7, 0x4D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)  # rectangle fond
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size, color=WHITE, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=8, bullet=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for i, item in enumerate(runs):
        txt, col, bd = (item + (False,))[:3] if isinstance(item, tuple) else (item, color, bold)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = col
        r.font.bold = bd; r.font.italic = italic
        r.font.name = "Calibri"
    if bullet:
        p.text = ""  # reset then rewrite with bullet char
    return p


def bullets(tf, items, size=18, color=WHITE, space=10):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"


def title_bar(s, text, kicker=None):
    if kicker:
        tf = box(s, 0.7, 0.45, 12, 0.5)
        para(tf, kicker.upper(), 13, ACCENT, bold=True, first=True)
    tf = box(s, 0.7, 0.85, 12, 1.0)
    para(tf, text, 30, WHITE, bold=True, first=True)
    # filet accent
    ln = s.shapes.add_shape(1, Inches(0.75), Inches(1.75), Inches(1.4), Pt(4))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    ln.shadow.inherit = False


# ============ SLIDE 1 — TITRE ============
s = slide()
tf = box(s, 0.9, 2.1, 11.5, 2.6)
para(tf, "ACCÈS AUX SOINS EN ÎLE-DE-FRANCE", 40, WHITE, bold=True, first=True, space_after=4)
para(tf, "Où agir en priorité ?", 26, ACCENT, bold=True, space_after=20)
para(tf, "Un tableau de bord d'aide à la décision pour les élus locaux", 18, GREY)
tf = box(s, 0.9, 6.2, 11.5, 1.0)
para(tf, "SAKOA Etia-Anaëlle — Ynov — Data Visualisation", 14, GREY, first=True)
para(tf, "Commanditaire : maires & présidents d'EPCI franciliens", 14, GREY)

# ============ SLIDE 2 — LE PROBLÈME ============
s = slide(); title_bar(s, "Le problème", "Contexte")
tf = box(s, 0.75, 2.1, 7.2, 4.6)
bullets(tf, [
    "L'accès aux médecins se dégrade, même en Île-de-France.",
    "Les écarts sont énormes d'une commune à l'autre, parfois entre voisines.",
    "Les élus (maires, présidents d'EPCI) sont en première ligne pour agir : MSP, aides à l'installation, locaux…",
    "Mais ils manquent d'un outil simple pour savoir où le besoin est le plus fort.",
], size=19, space=16)
tf = box(s, 8.4, 2.3, 4.2, 4.0)
para(tf, "62 %", 54, RED, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "des communes franciliennes sous le seuil de sous-dotation en généralistes", 16, GREY, align=PP_ALIGN.CENTER)

# ============ SLIDE 3 — LA QUESTION ============
s = slide(); title_bar(s, "Deux questions, deux niveaux de décision", "Objectif")
tf = box(s, 0.9, 2.4, 5.5, 3.6)
para(tf, "OÙ AGIR ?", 24, ACCENT, bold=True, first=True, space_after=10)
para(tf, "Vue président d'EPCI :", 17, WHITE, bold=True, space_after=4)
para(tf, "comparer les communes du territoire et repérer les plus tendues.", 17, GREY)
tf = box(s, 6.9, 2.4, 5.5, 3.6)
para(tf, "SUR QUOI AGIR ?", 24, ORANGE, bold=True, first=True, space_after=10)
para(tf, "Vue maire :", 17, WHITE, bold=True, space_after=4)
para(tf, "voir, pour une commune, ce qui décroche vraiment par rapport à la moyenne régionale.", 17, GREY)

# ============ SLIDE 4 — LES DONNÉES ============
s = slide(); title_bar(s, "Les données", "Sources ouvertes")
tf = box(s, 0.75, 2.05, 11.8, 4.8)
bullets(tf, [
    "DREES — APL généralistes 2023 (accessibilité par habitant, niveau commune).",
    "DREES / RPPS — densité de 4 spécialistes pour 100 000 habitants (cardio, dermato, ophtalmo, gynéco).",
    "DREES — âge moyen des médecins 2025 (par département).",
    "INSEE (COG) — périmètre des 1 266 communes d'Île-de-France.",
    "BANATIC — rattachement commune → EPCI + population municipale.",
], size=19, space=14)
tf = box(s, 0.75, 6.35, 11.8, 0.8)
para(tf, "Toutes croisées par code INSEE. Périmètre : 1 266 communes franciliennes.",
     16, GREY, first=True, italic=True)

# ============ SLIDE 5 — MÉTHODE / SCORE ============
s = slide(); title_bar(s, "Le score de tension d'accès", "Méthode")
tf = box(s, 0.75, 2.0, 11.8, 1.2)
para(tf, "Un indice unique de 0 à 100 par commune. Plus il est haut, plus l'accès est tendu.",
     19, WHITE, first=True)
tf = box(s, 0.75, 3.15, 11.8, 3.0)
bullets(tf, [
    "3 dimensions à parts égales (1/3 chacune) : généralistes (APL), spécialistes (densité), âge des médecins.",
    "Chaque indicateur ramené sur 0–100 (normalisation min-max, bornée aux 2e/98e centiles pour ignorer les extrêmes).",
    "Indicateurs d'offre inversés (offre faible = tension haute) ; âge en direct (médecins âgés = renouvellement à risque).",
    "Résultat : un classement comparable, et un rang relatif au sein de chaque EPCI.",
], size=18, space=14)
tf = box(s, 0.75, 6.45, 11.8, 0.7)
para(tf, "Score observé sur l'IDF : de 17 (mieux doté) à 93 (le plus tendu).",
     16, ACCENT, first=True, italic=True)

# ============ SLIDE 6 — CONSTATS ============
s = slide(); title_bar(s, "Trois constats qui parlent aux élus", "Résultats")
cards = [
    ("62 %", "des communes sous le seuil\nde sous-dotation en généralistes", RED),
    ("× 13", "d'écart sur les dermatologues\nentre Paris et la Seine-Saint-Denis", ORANGE),
    ("52,8 ans", "âge moyen des médecins :\nun mur de départs en retraite", ACCENT),
]
x = 0.7
for big, txt, col in cards:
    tf = box(s, x, 2.7, 3.9, 3.2)
    para(tf, big, 44, col, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=8)
    for i, line in enumerate(txt.split("\n")):
        para(tf, line, 16, GREY, align=PP_ALIGN.CENTER, space_after=2)
    x += 4.15

# ============ SLIDE 7 — LE DASHBOARD ============
s = slide(); title_bar(s, "Le tableau de bord", "Livrable")
tf = box(s, 0.75, 2.0, 11.8, 1.0)
para(tf, "Publié sur Tableau Public, accessible depuis un navigateur, sans installation.",
     19, WHITE, first=True)
tf = box(s, 0.9, 3.2, 5.5, 3.4)
para(tf, "1 · Carte EPCI", 22, ACCENT, bold=True, first=True, space_after=8)
bullets(tf, [
    "Toute l'IDF colorée par tension.",
    "Filtre par intercommunalité.",
    "Rouge = prioritaire, vert = bien doté.",
], size=17, space=10)
tf = box(s, 6.9, 3.2, 5.5, 3.4)
para(tf, "2 · Profil commune", 22, ORANGE, bold=True, first=True, space_after=8)
bullets(tf, [
    "Une commune choisie.",
    "Ses indicateurs vs moyenne IDF.",
    "On voit ce qui décroche.",
], size=17, space=10)

# ============ SLIDE 8 — DÉMO CARTE ============
s = slide(); title_bar(s, "Où agir : la carte EPCI", "Démo 1")
tf = box(s, 0.75, 2.0, 11.8, 1.4)
para(tf, "Le président d'EPCI filtre son territoire et repère d'un coup d'œil les communes les plus tendues.",
     19, WHITE, first=True)
# zone capture
ph = s.shapes.add_shape(1, Inches(2.3), Inches(3.4), Inches(8.7), Inches(3.5))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
ph.line.color.rgb = GREY; ph.shadow.inherit = False
ph.text_frame.text = "▶  Insérer ici la capture « Carte EPCI » du dashboard"
ph.text_frame.paragraphs[0].font.color.rgb = GREY
ph.text_frame.paragraphs[0].font.size = Pt(16)
ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============ SLIDE 9 — DÉMO PROFIL ============
s = slide(); title_bar(s, "Sur quoi agir : le profil commune", "Démo 2")
tf = box(s, 0.75, 2.0, 11.8, 1.4)
para(tf, "Pour la commune choisie : ses valeurs (orange) face à la moyenne régionale (vert), indicateur par indicateur.",
     19, WHITE, first=True)
ph = s.shapes.add_shape(1, Inches(2.3), Inches(3.4), Inches(8.7), Inches(3.5))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
ph.line.color.rgb = GREY; ph.shadow.inherit = False
ph.text_frame.text = "▶  Insérer ici la capture « Ma commune vs moyenne IDF »"
ph.text_frame.paragraphs[0].font.color.rgb = GREY
ph.text_frame.paragraphs[0].font.size = Pt(16)
ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============ SLIDE 10 — LIMITES ============
s = slide(); title_bar(s, "Limites à garder en tête", "Honnêteté")
tf = box(s, 0.75, 2.1, 11.8, 4.6)
bullets(tf, [
    "Densité des spécialistes disponible au département, pas à la commune : à lire comme une tendance.",
    "L'APL date de 2023 ; l'offre bouge vite.",
    "Le score est un outil d'alerte et de priorisation, pas un diagnostic médical complet.",
    "Les 3 dimensions sont pondérées à parts égales : un choix discutable, ajustable selon les priorités de l'élu.",
], size=19, space=16)

# ============ SLIDE 11 — LEVIERS ============
s = slide(); title_bar(s, "Des constats aux leviers d'action", "Actionnable")
tf = box(s, 0.75, 2.05, 11.8, 4.8)
bullets(tf, [
    "Maisons de santé pluriprofessionnelles (MSP) sur les communes les plus rouges.",
    "Aides à l'installation et médecins salariés là où l'offre libérale décroche.",
    "Télémédecine pour les spécialités les plus rares (dermato, ophtalmo).",
    "Coopération à l'échelle de l'EPCI : mutualiser plutôt que se concurrencer entre communes.",
    "Anticiper les départs en retraite là où l'âge moyen est le plus élevé.",
], size=19, space=14)

# ============ SLIDE 12 — CONCLUSION ============
s = slide()
tf = box(s, 0.9, 2.3, 11.5, 2.4)
para(tf, "Un outil simple pour une décision locale", 32, WHITE, bold=True, first=True, space_after=16)
para(tf, "Où agir → la carte. Sur quoi agir → le profil commune.", 20, ACCENT)
tf = box(s, 0.9, 5.0, 11.5, 1.6)
para(tf, "Dashboard en ligne :", 15, GREY, first=True, space_after=2)
para(tf, "public.tableau.com/app/profile/etia.ana.lle.sakoa", 16, ACCENT)
para(tf, "Merci — questions ?", 18, WHITE, bold=True, space_after=0)

out = "project_dashboard_health_access/deliverables/support_presentation.pptx"
prs.save(out)
print("OK ->", out, "|", len(prs.slides._sldIdLst), "slides")
